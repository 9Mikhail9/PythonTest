import os
import pandas as pd
from tkinter import *
from pptx import Presentation
import docx
from PyPDF2 import PdfFileReader

REPORT = "/Report.xlsx"
REPORTNUM = 0
symbols_list = ["\n","Unnamed:","NaN",",",".","?","!","(",")",";",":"," ’","’ ","'","‘ "," ‘","— "," —","1","2","3","4","5","6","7","8","9","0",'"',"\\","/"]

def open_catalogname(catalogname):
    if os.path.exists(catalogname.replace("\\", "/")):
        return True
    else:
        return False

def read_txt(filename):
    with open(filename, "r") as file:
        text = file.read()
    return text


def read_pdf(filename):
    text = []
    pdf = PdfFileReader(filename)
    for num_page in range(len(pdf.pages)):
        page = pdf.pages[num_page]
        text.append(page.extract_text())
    return str(text)


def read_docx(filename):
    doc = docx.Document(filename)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return str(text)


def read_xlsx(filename):
    text = pd.read_excel(filename, index_col=0)
    return str(text)

def read_pptx(filename):
    ppt = Presentation(filename)
    text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text.append(run.text)
    return str(text)

def get_words(text):
    for symbol in range(len(symbols_list)):
        text = text.replace(symbols_list[symbol], " ")
    text = text.lower()
    words = text.split()
    words.sort()
    return words




def get_words_dict(words):
    words_dict = dict()

    for word in words:
        if word in words_dict:
            words_dict[word] = words_dict[word] + 1
        else:
            words_dict[word] = 1
    return words_dict


def get_file_name(FilePath, chk_txt, chk_doc, chk_pdf, chk_pptx, chk_xlsx):
    all_files = []
    for file in os.listdir(FilePath):
        if file.endswith(".txt") and chk_txt:
            all_files.append(file)
        if ((file.endswith(".doc") or file.endswith(".docx")) and chk_doc):
            all_files.append(file)
        if file.endswith(".pdf") and chk_pdf:
            all_files.append(file)
        if ((file.endswith(".pptx") or file.endswith(".ppt")) and chk_pptx):
            all_files.append(file)
        if ((file.endswith(".xlsx") or file.endswith(".xls")) and chk_xlsx):
            all_files.append(file)
    return all_files


def opening_file(filename, catalogname):
    if not os.path.exists((catalogname + '\ ' + filename).replace(" ", "")):
        return False
    else:
        return True


def create_files_tatus_table(data_file, filename, status, mylist):
    if status:
        new_row = pd.DataFrame([[filename, 'Read']], columns=['File Name', 'Read/Not Read'])
        status = filename + ":           Read"
    else:
        new_row = pd.DataFrame([[filename, 'Not Read']], columns=['File Name', 'Read/Not Read'])
        status = filename + ":           Not Read"
    if len(status) >= 55:
        status = status[:5] + "..." + status[len(status) - 47:]
        mylist.insert(END, status)
    else:
        mylist.insert(END, status)
    data_file = pd.concat([data_file, new_row], ignore_index=True)
    return data_file


def create_word_quantity_table(data_word, word, quantity):
    new_row = pd.DataFrame([[word, quantity]], columns=['Word', 'Quantity'])
    data_word = pd.concat([data_word, new_row], ignore_index=True)
    return data_word


def saving_report(data_word, data_file, catalogname):
    global REPORT, REPORTNUM
    if os.path.exists(catalogname.replace("\\", "/") + "\Report Catalog"):
        if not os.path.exists(catalogname + "\Report Catalog" + "\\" + REPORT):
            writer = pd.ExcelWriter(catalogname + "\Report Catalog" + "\\" + REPORT, engine='xlsxwriter')
            data_file.to_excel(writer, 'Sheet1', index_label=False, index=False, header=True)
            data_word.to_excel(writer, 'Sheet2', index_label=False, index=False, header=True)
            writer.save()
        else:
            REPORTNUM = REPORTNUM + 1
            REPORT = '/Report' + str(REPORTNUM) + '.xlsx'
            saving_report(data_word, data_file, catalogname)
    else:
        os.mkdir(catalogname + "\Report Catalog")
        saving_report(data_word, data_file, catalogname)
    data_file.reset_index()
